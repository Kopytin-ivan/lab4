from __future__ import annotations

from pathlib import Path

import pandas as pd
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = PROJECT_ROOT / "reports"
FINAL_DIR = REPORTS_DIR / "final"
FIGURES_DIR = REPORTS_DIR / "figures"
MODELS_DIR = REPORTS_DIR / "models"
ANALYTICS_DIR = REPORTS_DIR / "analytics"
PRESENTATION_DIR = PROJECT_ROOT / "presentation"


def fmt(value: float, digits: int = 3) -> str:
    return f"{float(value):.{digits}f}"


def to_markdown_table(df: pd.DataFrame) -> str:
    columns = list(df.columns)
    lines = [
        "| " + " | ".join(columns) + " |",
        "| " + " | ".join(["---"] * len(columns)) + " |",
    ]
    for _, row in df.iterrows():
        values = []
        for col in columns:
            value = row[col]
            if isinstance(value, float):
                value = fmt(value)
            values.append(str(value))
        lines.append("| " + " | ".join(values) + " |")
    return "\n".join(lines)


def load_results() -> dict[str, pd.DataFrame]:
    return {
        "metrics": pd.read_csv(MODELS_DIR / "model_metrics.csv"),
        "cv": pd.read_csv(MODELS_DIR / "cv_metrics.csv"),
        "thresholds": pd.read_csv(ANALYTICS_DIR / "quality_thresholds.csv"),
        "summary": pd.read_csv(ANALYTICS_DIR / "condition_prediction_summary.csv"),
        "rec32": pd.read_csv(ANALYTICS_DIR / "recommendations_ra_3_2.csv"),
        "rec20": pd.read_csv(ANALYTICS_DIR / "recommendations_ra_2_0.csv"),
    }


def add_doc_table(document: Document, df: pd.DataFrame, columns: list[str], max_rows: int = 8) -> None:
    shown = df[columns].head(max_rows)
    table = document.add_table(rows=1, cols=len(columns))
    table.style = "Table Grid"
    for index, col in enumerate(columns):
        table.rows[0].cells[index].text = col
    for _, row in shown.iterrows():
        cells = table.add_row().cells
        for index, col in enumerate(columns):
            value = row[col]
            if isinstance(value, float):
                value = fmt(value)
            cells[index].text = str(value)


def create_markdown_report(results: dict[str, pd.DataFrame]) -> Path:
    metrics = results["metrics"]
    cv = results["cv"]
    thresholds = results["thresholds"]
    rec32 = results["rec32"]
    rec20 = results["rec20"]

    best = metrics.iloc[0]
    best_threshold = thresholds.iloc[0]
    best_rec32 = rec32.iloc[0]
    best_rec20 = rec20.iloc[0]

    report = f"""# Отчет по проекту

## Тема

Предсказание шероховатости обработанной поверхности по шуму во время фрезерования.

## Цель

Разработать функциональный прототип аналитического приложения, которое по параметрам режима обработки и акустическим признакам оценивает шероховатость поверхности `Ra`, риск выхода за допуск и рекомендует режимы обработки.

## Данные

Использован датасет Milling Surface Roughness Acoustic Sensor Dataset. В проекте использованы Audacity-проекты `.aup` с аудио-блоками `.au` и Excel-таблица `SurfaceRoughnessValuesConsolidated.xlsx` со значениями шероховатости.

После подготовки получено:

- 319 аудиозаписей;
- 16 уникальных режимов обработки;
- 42 столбца в итоговой таблице признаков;
- целевая переменная: `roughness_ra`.

## Методика

1. Подготовлены метаданные: подача, скорость шпинделя, глубина резания, номер запуска и метка `Ra`.
2. Из аудио извлечены статистические и спектральные признаки: RMS, crest factor, zero crossing rate, спектральный центр, rolloff, энергия в частотных диапазонах.
3. Обучены модели регрессии: Linear Regression, Random Forest, Gradient Boosting.
4. Оценка качества выполнена с разбиением по `condition_id`, чтобы одинаковый режим не попадал одновременно в train и test.
5. Выполнена аналитика порогов и рекомендаций режимов.

## Результаты моделей

Лучшая модель на отложенных режимах: **{best['model']}**.

- MAE: {fmt(best['test_MAE'])}
- RMSE: {fmt(best['test_RMSE'])}
- R2: {fmt(best['test_R2'])}

Кросс-валидация GroupKFold:

{to_markdown_table(cv)}

## Порог качества

В качестве допуска принят уровень `Ra <= 3.2`. Лучший акустический признак для отделения плохого качества:

- признак: `{best_threshold['feature']}`;
- правило: `{best_threshold['feature']} {best_threshold['direction']} {fmt(best_threshold['threshold'], 6)}`;
- доля плохого качества после порога: {fmt(best_threshold['bad_rate_after_threshold'])};
- доля плохого качества до порога: {fmt(best_threshold['bad_rate_before_threshold'])}.

Интерпретация: при росте акустической энергии в диапазоне 2000-8000 Hz выше найденного порога риск выхода шероховатости за допуск заметно возрастает.

## Рекомендации режимов

Для допуска `Ra <= 3.2` лучший режим по производительности среди годных:

- подача: {best_rec32['feed_mm']} мм;
- скорость: {int(best_rec32['speed_rpm'])} rpm;
- глубина: {best_rec32['depth_mm']} мм;
- фактическая Ra: {fmt(best_rec32['roughness_ra'])};
- прогноз Ra: {fmt(best_rec32['predicted_ra_mean'])}.

Для более чистой обработки `Ra <= 2.0` лучший режим:

- подача: {best_rec20['feed_mm']} мм;
- скорость: {int(best_rec20['speed_rpm'])} rpm;
- глубина: {best_rec20['depth_mm']} мм;
- фактическая Ra: {fmt(best_rec20['roughness_ra'])};
- прогноз Ra: {fmt(best_rec20['predicted_ra_mean'])}.

## Функциональный прототип

Создан Streamlit-прототип `app/streamlit_app.py`. Он позволяет:

- выбрать режим обработки;
- изменить акустические признаки;
- получить прогноз `Ra`;
- увидеть предупреждение о превышении акустического порога;
- подобрать режим под целевую шероховатость.

Команда запуска:

```powershell
.\\.venv\\Scripts\\streamlit.exe run app\\streamlit_app.py
```

## Вывод

Проект реализует data-driven подход: данные производственного процесса превращены в признаки, обучена модель машинного обучения, построены интерпретируемые пороги риска и рекомендации режимов обработки. На малом числе уникальных режимов линейная модель оказалась устойчивее ансамблевых моделей, которые сильнее переобучались.
"""

    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    path = FINAL_DIR / "surface_roughness_report.md"
    path.write_text(report, encoding="utf-8")
    return path


def create_docx_report(results: dict[str, pd.DataFrame]) -> Path:
    metrics = results["metrics"]
    cv = results["cv"]
    thresholds = results["thresholds"]
    rec32 = results["rec32"]
    rec20 = results["rec20"]

    best = metrics.iloc[0]
    best_threshold = thresholds.iloc[0]

    document = Document()
    document.add_heading("Отчет по проекту", level=0)
    document.add_paragraph("Тема: предсказание шероховатости обработанной поверхности по шуму во время фрезерования.")

    document.add_heading("Цель работы", level=1)
    document.add_paragraph(
        "Разработать функциональный прототип приложения, которое использует data-driven подход "
        "и машинное обучение для прогноза шероховатости поверхности Ra."
    )

    document.add_heading("Данные", level=1)
    document.add_paragraph(
        "Использован Milling Surface Roughness Acoustic Sensor Dataset: аудиозаписи фрезерования "
        "и таблица измерений шероховатости. После обработки получено 319 записей и 16 режимов обработки."
    )

    document.add_heading("Методика", level=1)
    for item in [
        "Подготовка метаданных и связывание аудиозаписей с меткой Ra.",
        "Извлечение статистических и спектральных признаков из акустического сигнала.",
        "Обучение моделей Linear Regression, Random Forest, Gradient Boosting.",
        "Оценка качества с разбиением по режимам обработки condition_id.",
        "Поиск порога акустического признака и подбор оптимальных режимов.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("Результаты моделирования", level=1)
    document.add_paragraph(
        f"Лучшая модель: {best['model']}. "
        f"Test MAE = {fmt(best['test_MAE'])}, RMSE = {fmt(best['test_RMSE'])}, R2 = {fmt(best['test_R2'])}."
    )
    add_doc_table(document, metrics, ["model", "test_MAE", "test_RMSE", "test_R2"], max_rows=3)

    document.add_paragraph("Результаты GroupKFold-кросс-валидации:")
    add_doc_table(document, cv, ["model", "cv_MAE_mean", "cv_RMSE_mean", "cv_R2_mean"], max_rows=3)

    document.add_heading("Порог качества", level=1)
    document.add_paragraph(
        f"Для допуска Ra <= 3.2 лучший порог найден по признаку {best_threshold['feature']}: "
        f"{best_threshold['feature']} {best_threshold['direction']} {fmt(best_threshold['threshold'], 6)}. "
        f"После порога доля плохого качества равна {fmt(best_threshold['bad_rate_after_threshold'])}, "
        f"до порога - {fmt(best_threshold['bad_rate_before_threshold'])}."
    )
    if (FIGURES_DIR / "quality_threshold_best_feature.png").exists():
        document.add_picture(str(FIGURES_DIR / "quality_threshold_best_feature.png"), width=Inches(5.7))

    document.add_heading("Рекомендации режимов", level=1)
    document.add_paragraph("Лучшие режимы для Ra <= 3.2:")
    add_doc_table(
        document,
        rec32,
        ["condition_id", "feed_mm", "speed_rpm", "depth_mm", "roughness_ra", "predicted_ra_mean", "productivity_proxy"],
        max_rows=6,
    )
    document.add_paragraph("Лучшие режимы для более чистой обработки Ra <= 2.0:")
    add_doc_table(
        document,
        rec20,
        ["condition_id", "feed_mm", "speed_rpm", "depth_mm", "roughness_ra", "predicted_ra_mean", "productivity_proxy"],
        max_rows=6,
    )

    document.add_heading("Функциональный прототип", level=1)
    document.add_paragraph(
        "Создан Streamlit-прототип app/streamlit_app.py. Приложение прогнозирует Ra, "
        "показывает риск выхода за допуск и рекомендует режимы обработки."
    )
    document.add_paragraph(r".\.venv\Scripts\streamlit.exe run app\streamlit_app.py")

    document.add_heading("Вывод", level=1)
    document.add_paragraph(
        "Проект демонстрирует применение машинного обучения к данным производственного процесса. "
        "Получены модель прогноза, порог акустического риска и рекомендации режимов обработки."
    )

    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    path = FINAL_DIR / "surface_roughness_report.docx"
    document.save(path)
    return path


def set_title(slide, title: str, subtitle: str | None = None) -> None:
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(30)
    title_shape.text_frame.paragraphs[0].font.bold = True
    if subtitle:
        box = slide.shapes.add_textbox(PptInches(0.75), PptInches(1.55), PptInches(8.8), PptInches(0.6))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(80, 80, 80)


def add_bullets(slide, bullets: list[str], left=0.85, top=1.55, width=8.8, height=4.8) -> None:
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(18)


def add_picture_if_exists(slide, name: str, left=5.35, top=1.55, width=4.2) -> None:
    path = FIGURES_DIR / name
    if path.exists():
        slide.shapes.add_picture(str(path), PptInches(left), PptInches(top), width=PptInches(width))


def create_presentation(results: dict[str, pd.DataFrame]) -> Path:
    metrics = results["metrics"]
    thresholds = results["thresholds"]
    rec32 = results["rec32"]
    best = metrics.iloc[0]
    best_threshold = thresholds.iloc[0]
    best_rec = rec32.iloc[0]

    presentation = Presentation()
    presentation.slide_width = PptInches(10)
    presentation.slide_height = PptInches(5.625)

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Прогноз шероховатости поверхности"
    slide.placeholders[1].text = "Data-driven анализ акустики фрезерования и ML-прототип"

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Постановка задачи")
    add_bullets(
        slide,
        [
            "Цель: предсказать Ra по шуму, подаче, скорости шпинделя и глубине резания.",
            "Данные: Milling Surface Roughness Acoustic Sensor Dataset.",
            "Результат: модель, порог риска, рекомендации режимов и Streamlit-прототип.",
        ],
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Подготовка данных")
    add_bullets(
        slide,
        [
            "319 аудиозаписей Audacity .aup/.au.",
            "16 уникальных режимов фрезерования.",
            "Целевая переменная: roughness_ra.",
            "Связь аудио с метками выполнена через feed/speed/doc в структуре папок.",
        ],
        width=4.3,
    )
    add_picture_if_exists(slide, "runs_per_condition.png")

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Извлечение признаков")
    add_bullets(
        slide,
        [
            "Статистика сигнала: RMS, peak-to-peak, crest factor.",
            "Спектр: centroid, bandwidth, rolloff, энергия в частотных диапазонах.",
            "Технологические признаки: feed, speed, depth.",
        ],
        width=4.3,
    )
    add_picture_if_exists(slide, "feature_correlations.png")

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Сравнение моделей")
    add_bullets(
        slide,
        [
            f"Лучшая модель: {best['model']}.",
            f"Test MAE: {fmt(best['test_MAE'])}.",
            f"Test RMSE: {fmt(best['test_RMSE'])}.",
            f"Test R2: {fmt(best['test_R2'])}.",
            "Деление train/test выполнено по condition_id.",
        ],
        width=4.3,
    )
    add_picture_if_exists(slide, "model_cv_mae.png")

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Качество прогноза")
    add_bullets(
        slide,
        [
            "Линейная регрессия лучше обобщилась на новые режимы.",
            "Random Forest и Gradient Boosting переобучались из-за малого числа режимов.",
            "Модель пригодна как прототип поддержки технолога.",
        ],
        width=4.3,
    )
    add_picture_if_exists(slide, "best_model_actual_vs_predicted.png")

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Порог ухудшения качества")
    add_bullets(
        slide,
        [
            "Допуск принят: Ra <= 3.2.",
            f"Лучший порог: {best_threshold['feature']} {best_threshold['direction']} {fmt(best_threshold['threshold'], 6)}.",
            f"Доля плохого качества после порога: {fmt(best_threshold['bad_rate_after_threshold'])}.",
            "Рост энергии 2000-8000 Hz связан с ухудшением поверхности.",
        ],
        width=4.3,
    )
    add_picture_if_exists(slide, "quality_threshold_best_feature.png")

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Оптимальные режимы")
    add_bullets(
        slide,
        [
            "Режимы ранжируются по прогнозу Ra и proxy производительности.",
            f"Лучший для Ra <= 3.2: feed {best_rec['feed_mm']}, speed {int(best_rec['speed_rpm'])}, depth {best_rec['depth_mm']}.",
            f"Прогноз Ra: {fmt(best_rec['predicted_ra_mean'])}.",
        ],
        width=4.3,
    )
    add_picture_if_exists(slide, "recommended_modes_ra_3_2.png")

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Функциональный прототип")
    add_bullets(
        slide,
        [
            "Streamlit-приложение app/streamlit_app.py.",
            "Ввод режима и акустических признаков.",
            "Прогноз Ra и статус годности.",
            "Рекомендации режимов под заданную чистоту.",
        ],
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    set_title(slide, "Выводы")
    add_bullets(
        slide,
        [
            "Проект реализует полный data-driven цикл: данные -> признаки -> модель -> аналитика -> прототип.",
            "Порог акустической энергии помогает оценивать риск плохого качества.",
            "При скорости 500 rpm чаще достигается лучшая чистота поверхности.",
            "Прототип можно развивать: добавить загрузку нового аудио и дообучение модели.",
        ],
    )

    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    path = PRESENTATION_DIR / "surface_roughness_presentation.pptx"
    presentation.save(path)
    return path


def main() -> None:
    results = load_results()
    md_path = create_markdown_report(results)
    docx_path = create_docx_report(results)
    pptx_path = create_presentation(results)

    print("Created:")
    print(md_path)
    print(docx_path)
    print(pptx_path)


if __name__ == "__main__":
    main()
